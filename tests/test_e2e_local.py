from __future__ import annotations

import hashlib
import os
from pathlib import Path

import pytest

from ai_material_preprocessor.converters.common import run_command
from ai_material_preprocessor.converters.markdown import to_markdown
from ai_material_preprocessor.converters.office_pdf import to_pdf
from ai_material_preprocessor.converters.video import (
    compress,
    extract_audio,
    keyframes_contact_sheet,
    rename_copy,
    standardize,
)
from ai_material_preprocessor.services.config import load_config
from ai_material_preprocessor.services.document_enhancement import EnhancementOptions
from ai_material_preprocessor.services.environment import detect_tools
from ai_material_preprocessor.services.ocr import RapidOCREngine


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


@pytest.fixture(scope="module")
def detected_tools():
    return detect_tools(load_config())


def test_real_markitdown_html_conversion(tmp_path: Path) -> None:
    source = Path(__file__).parent / "fixtures" / "sample.html"
    result = to_markdown(source, tmp_path)
    content = result.read_text(encoding="utf-8")
    assert "AI 素材预处理" in content
    assert "最小验证文件" in content


def test_real_markitdown_pptx_conversion(tmp_path: Path) -> None:
    from pptx import Presentation

    source = tmp_path / "lesson.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "数字逻辑课程"
    slide.placeholders[1].text = "组合逻辑与时序逻辑"
    presentation.save(source)

    result = to_markdown(source, tmp_path / "outputs")
    content = result.read_text(encoding="utf-8")
    assert "数字逻辑课程" in content
    assert "组合逻辑与时序逻辑" in content


def test_real_markitdown_xlsx_conversion(tmp_path: Path) -> None:
    from openpyxl import Workbook

    source = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "成绩"
    sheet.append(["姓名", "分数"])
    sheet.append(["小明", 95])
    workbook.save(source)

    result = to_markdown(source, tmp_path / "outputs")
    content = result.read_text(encoding="utf-8")
    assert "姓名" in content
    assert "小明" in content


def test_real_enhanced_pptx_and_xlsx_structure(tmp_path: Path) -> None:
    from openpyxl import Workbook
    from pptx import Presentation

    presentation = Presentation()
    for index in range(1, 4):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Slide {index}"
        slide.placeholders[1].text = f"Repeated Template\nUnique {index}"
    pptx = tmp_path / "deck.pptx"
    presentation.save(pptx)
    ppt_result = to_markdown(
        pptx,
        tmp_path / "outputs",
        enhance=True,
        enhancement_options=EnhancementOptions(split_enabled=True, target_tokens=40, max_tokens=80),
    )
    ppt_content = ppt_result.read_text(encoding="utf-8")
    assert ppt_content.count("Repeated Template") == 0
    assert ppt_content.count("\n---\n") == 2
    assert (ppt_result.parent / "raw.md").is_file()
    assert (ppt_result.parent / "manifest.json").is_file()

    workbook = Workbook()
    workbook.active.title = "Scores"
    workbook.active.append(["Name", "Score"])
    workbook.active.append(["Alice", 95])
    workbook.create_sheet("Stats").append(["Metric", "Value"])
    xlsx = tmp_path / "book.xlsx"
    workbook.save(xlsx)
    xlsx_result = to_markdown(
        xlsx,
        tmp_path / "outputs",
        enhance=True,
        enhancement_options=EnhancementOptions(split_enabled=False),
    )
    xlsx_content = xlsx_result.read_text(encoding="utf-8")
    assert "# 工作表：Scores" in xlsx_content
    assert "# 工作表：Stats" in xlsx_content


def test_real_local_ocr_image(tmp_path: Path) -> None:
    from PIL import Image, ImageDraw, ImageFont

    font_path = Path(r"C:\Windows\Fonts\arial.ttf")
    if not font_path.is_file():
        pytest.skip("Arial font is not available")
    source = tmp_path / "ocr.png"
    image = Image.new("RGB", (1000, 240), "white")
    ImageDraw.Draw(image).text(
        (40, 50), "HELLO 123", fill="black", font=ImageFont.truetype(str(font_path), 96)
    )
    image.save(source)

    result = RapidOCREngine().extract(source)

    assert result
    recognized = " ".join(item[1] for item in result)
    assert "HELLO" in recognized and "123" in recognized


def test_real_ffmpeg_video_pipeline_preserves_source(tmp_path: Path, detected_tools) -> None:
    ffmpeg = detected_tools["ffmpeg"].path
    if not ffmpeg:
        pytest.skip("FFmpeg is not available")
    source = tmp_path / "source.mp4"
    run_command(
        [
            ffmpeg,
            "-hide_banner",
            "-loglevel",
            "error",
            "-y",
            "-f",
            "lavfi",
            "-i",
            "testsrc2=size=320x240:rate=25",
            "-f",
            "lavfi",
            "-i",
            "sine=frequency=440:sample_rate=44100",
            "-t",
            "1",
            "-shortest",
            "-metadata",
            "creation_time=2026-07-31T07:30:21Z",
            "-metadata",
            "location=+30.2512+120.1693/",
            "-c:v",
            "libx264",
            "-pix_fmt",
            "yuv420p",
            "-c:a",
            "aac",
            str(source),
        ]
    )
    original_hash = _sha256(source)
    output_root = tmp_path / "outputs"

    compressed = compress(source, output_root, ffmpeg, 28, "veryfast")
    audio = extract_audio(source, output_root, ffmpeg, "mp3", "128k")
    standardized = standardize(source, output_root, ffmpeg)
    contact_sheet = keyframes_contact_sheet(
        source, output_root, ffmpeg, scene_threshold=0.2, max_frames=8, columns=4
    )
    renamed = rename_copy(
        source,
        output_root,
        detected_tools["ffprobe"].path,
        "杭州西湖",
        exiftool=detected_tools["exiftool"].path,
        ffmpeg=ffmpeg,
    )

    assert all(
        path.is_file() and path.stat().st_size > 0
        for path in (
            compressed,
            audio,
            standardized,
            contact_sheet,
            renamed,
        )
    )
    assert (contact_sheet.parent / "manifest.json").is_file()
    assert list((contact_sheet.parent / "frames").glob("frame_*.jpg"))
    assert "杭州西湖" in renamed.name
    assert _sha256(source) == original_hash


def test_real_office_pdf_when_document_is_provided(tmp_path: Path, detected_tools) -> None:
    configured = os.environ.get("AI_MATERIAL_E2E_DOCX", "")
    if not configured or not Path(configured).is_file():
        pytest.skip("Set AI_MATERIAL_E2E_DOCX to run the Microsoft Office PDF test")
    source = Path(configured)
    original_hash = _sha256(source)

    result = to_pdf(
        source,
        tmp_path,
        detected_tools["libreoffice"].path,
        detected_tools["winword"].path,
        detected_tools["powerpoint"].path,
    )

    assert result.is_file() and result.stat().st_size > 0
    assert result.read_bytes().startswith(b"%PDF")
    assert _sha256(source) == original_hash
    markdown = to_markdown(result, tmp_path / "markdown")
    assert markdown.is_file() and markdown.stat().st_size > 0
